"""
CloudMundi Trading Platform Analysis - PowerPoint Generator
Creates a professional PowerPoint presentation from the analysis data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import os

# CloudMundi colors
CLOUDMUNDI_TEAL = RGBColor(45, 212, 191)
CLOUDMUNDI_NAVY = RGBColor(30, 58, 138)
CLOUDMUNDI_DARK = RGBColor(15, 23, 42)

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9 aspect ratio
    prs.slide_height = Inches(9)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add header image if exists
    if os.path.exists('header_cloudmundi.png'):
        left = Inches(0)
        top = Inches(0)
        slide1.shapes.add_picture('header_cloudmundi.png', left, top, width=prs.slide_width)
    
    # Add title text
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(3))
    title_frame = title_box.text_frame
    title_frame.text = "Comparative Analysis: eToro vs Revolut vs Robinhood"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Executive Summary"
    
    content = slide2.placeholders[1].text_frame
    content.text = "This report evaluates three leading trading platforms to identify the most suitable option for diversified investment across four key asset classes:"
    
    p = content.add_paragraph()
    p.text = "• Bitcoin (BTC)"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• Commodities"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• S&P 500 Index"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "• Bonds"
    p.level = 1
    
    # Slide 3: Investment Scope
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Investment Scope"
    
    content = slide3.placeholders[1].text_frame
    content.text = "Bitcoin: High-growth digital asset with increasing institutional adoption"
    
    p = content.add_paragraph()
    p.text = "Commodities: Inflation hedge and portfolio diversifier (e.g., Gold)"
    
    p = content.add_paragraph()
    p.text = "S&P 500: Core equity market exposure to U.S. economy"
    
    p = content.add_paragraph()
    p.text = "Bonds: Capital preservation and yield generation via fixed-income exposure"
    
    # Slide 4: Platform Feature Matrix
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide4.shapes.title.text = "Platform Feature Comparison Matrix"
    
    if os.path.exists('platform_feature_matrix.png'):
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(6)
        slide4.shapes.add_picture('platform_feature_matrix.png', left, top, height=height)
    
    # Slide 5: Platform Fee Table
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "Platform Fee Comparison"
    
    if os.path.exists('platform_fee_table.png'):
        left = Inches(2)
        top = Inches(2.5)
        height = Inches(3.5)
        slide5.shapes.add_picture('platform_fee_table.png', left, top, height=height)
        
    # Slide 5a: Asset Predictions
    slide5a = prs.slides.add_slide(prs.slide_layouts[5])
    slide5a.shapes.title.text = "Asset Class Performance Predictions (1 Year)"
    
    if os.path.exists('asset_predictions.png'):
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(6)
        slide5a.shapes.add_picture('asset_predictions.png', left, top, height=height)
        
    # Slide 5b: Platform Fee Impact
    slide5b = prs.slides.add_slide(prs.slide_layouts[5])
    slide5b.shapes.title.text = "Platform Fee Impact Analysis"
    
    if os.path.exists('platform_fee_impact.png'):
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(6.5)
        slide5b.shapes.add_picture('platform_fee_impact.png', left, top, height=height)
    
    # Slide 6: Portfolio Projection - 1 Month
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    slide6.shapes.title.text = "Portfolio Projection - 1 Month (£10,000 Investment)"
    
    if os.path.exists('portfolio_projection_1month.png'):
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(6)
        slide6.shapes.add_picture('portfolio_projection_1month.png', left, top, height=height)
    
    # Slide 7: Portfolio Projection - 1 Year
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    slide7.shapes.title.text = "Portfolio Projection - 1 Year (£10,000 Investment)"
    
    if os.path.exists('portfolio_projection_1year.png'):
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(6)
        slide7.shapes.add_picture('portfolio_projection_1year.png', left, top, height=height)
    
    # Slide 8: Asset Allocation
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    slide8.shapes.title.text = "Recommended Asset Allocation"
    
    if os.path.exists('asset_allocation_pie.png'):
        left = Inches(3)
        top = Inches(1.5)
        height = Inches(6)
        slide8.shapes.add_picture('asset_allocation_pie.png', left, top, height=height)
    
    # Slide 9: Platform Comparison
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Platform Strengths & Limitations"
    
    content = slide9.placeholders[1].text_frame
    content.text = "eToro"
    content.paragraphs[0].font.bold = True
    
    p = content.add_paragraph()
    p.text = "✓ Multi-asset support, social trading, global coverage"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "✗ Higher trading spreads, limited U.S. access"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "\nRevolut"
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "✓ Mobile-first, integrated banking, beginner-friendly"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "✗ High crypto spreads, limited products"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "\nRobinhood"
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "✓ Commission-free, strong education, fractional shares"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "✗ U.S. only, no commodities, limited crypto wallet"
    p.level = 1
    
    # Slide 10: Recommendation
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "🏆 Professional Recommendation"
    
    content = slide10.placeholders[1].text_frame
    content.text = "Recommended Platform: eToro"
    content.paragraphs[0].font.bold = True
    content.paragraphs[0].font.size = Pt(24)
    
    p = content.add_paragraph()
    p.text = "\nRationale:"
    p.font.bold = True
    
    p = content.add_paragraph()
    p.text = "• Broadest functionality across all four investment categories"
    
    p = content.add_paragraph()
    p.text = "• Robust support for Bitcoin custody"
    
    p = content.add_paragraph()
    p.text = "• Access to commodities and diversified ETF offerings"
    
    p = content.add_paragraph()
    p.text = "• Global regulatory framework"
    
    p = content.add_paragraph()
    p.text = "\nNote: For U.S.-based users, Robinhood remains a strong alternative"
    p.font.italic = True
    
    # Save presentation
    prs.save('CloudMundi_Trading_Platform_Analysis.pptx')
    print("✅ PowerPoint presentation created: CloudMundi_Trading_Platform_Analysis.pptx")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("⚠️  python-pptx not installed. Install with: pip install python-pptx")
        print("📝 Manual PowerPoint creation guide:")
        print("\n1. Title Slide: Use header image + title")
        print("2. Executive Summary: Key points about the analysis")
        print("3. Investment Scope: 4 asset classes")
        print("4-8. Insert each chart as a full slide")
        print("9. Platform comparison table")
        print("10. Final recommendation") 